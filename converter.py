import os
import subprocess
import tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def convert_single_pptx(pptx_path, output_dir):
    """利用 LibreOffice 跨平台静默命令行将单个 PPTX 转换为 PDF 文件"""
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            user_profile = f"file://{temp_dir}"
            subprocess.run([
                "libreoffice",
                f"-env:UserInstallation={user_profile}",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                pptx_path
            ], check=True, capture_output=True, timeout=120)
        return pptx_path, True
    except Exception:
        return pptx_path, False

def pdf_to_png_high_speed(pdf_path, output_dir):
    """利用 PyMuPDF 高速渲染 PDF 为适合视觉大模型输入的高清 PNG 图片"""
    base_name = os.path.splitext(os.path.basename(pdf_path))[0]
    generated_images = []
    try:
        doc = fitz.open(pdf_path)
        for i, page in enumerate(doc):
            zoom = 200 / 72
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            png_name = f"{base_name}_slide_{i:04d}.png"
            png_path = os.path.join(output_dir, png_name)
            pix.save(png_path)
            generated_images.append(png_path)
        doc.close()
    except Exception as e:
        print(f"Error rendering {pdf_path} to image: {e}")
    return generated_images

def extract_text_from_shape(shape):
    """抓取 shape 中文字内容并将其转化带缩进的 markdown"""
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        level = para.level if hasattr(para, 'level') else 0
        indent = "  " * level

        para_text = ""
        for run in para.runs:
            text = run.text
            if not text.strip():
                para_text += text
                continue
            if run.font.bold:
                text = f"**{text}**"
            if run.font.italic:
                text = f"*{text}*"
            para_text += text

        if not para_text.strip():
            continue

        if para_text.startswith(('•', '-', '*', '●', '◆', '►', '□', '■', '✓', '→')):
            parts.append(f"{indent}{para_text}")
        elif level > 0:
            parts.append(f"{indent}- {para_text}")
        else:
            parts.append(f"{indent}{para_text}")
    return "\n".join(parts)

def extract_table_markdown(shape):
    """提取幻灯片内原生表格结构并构造其标准的 Markdown Table"""
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        return ""
    table = shape.table
    rows_md = []
    for row_idx, row in enumerate(table.rows):
        cells_md = []
        for cell in row.cells:
            cells_md.append(cell.text.replace("\n", " ").strip())
        rows_md.append("| " + " | ".join(cells_md) + " |")
        if row_idx == 0:
            rows_md.append("| " + " | ".join(["---"] * len(row.cells)) + " |")
    return "\n".join(rows_md)

def slide_to_markdown_layout_aware(slide, slide_width, slide_height):
    """根据空间几何关系，执行基于流式布局（通栏、分双栏）的对齐抽取"""
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    table_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
    all_shapes = text_shapes + table_shapes

    if not all_shapes:
        return ""

    full_width_threshold = slide_width * 0.7
    mid_x = slide_width / 2

    banners = []
    columns_left = []
    columns_right = []

    for shape in all_shapes:
        center_x = shape.left + shape.width / 2
        if shape.width >= full_width_threshold:
            banners.append(shape)
        elif center_x < mid_x:
            columns_left.append(shape)
        else:
            columns_right.append(shape)

    sorted_banners = sorted(banners, key=lambda s: s.top)
    sorted_left = sorted(columns_left, key=lambda s: s.top)
    sorted_right = sorted(columns_right, key=lambda s: s.top)

    top_banners = [s for s in sorted_banners if s.top < slide_height * 0.2]
    bottom_banners = [s for s in sorted_banners if s.top >= slide_height * 0.2]

    ordered_shapes = top_banners + sorted_left + sorted_right + bottom_banners

    blocks = []
    for shape in ordered_shapes:
        if shape.has_text_frame:
            text = extract_text_from_shape(shape)
            if text.strip():
                blocks.append(text)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_md = extract_table_markdown(shape)
            if table_md.strip():
                blocks.append(table_md)

    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            blocks.append(f"\n\n<!-- Speaker Notes:\n{notes_text}\n-->")

    return "\n\n".join(blocks)
